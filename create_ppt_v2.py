# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C1 = RGBColor(0x0B, 0x3D, 0x91)
C2 = RGBColor(0x15, 0x65, 0xC0)
C3 = RGBColor(0x00, 0x96, 0x88)
C4 = RGBColor(0xD8, 0x43, 0x35)
C5 = RGBColor(0xEF, 0x6C, 0x00)
CW = RGBColor(0xFF, 0xFF, 0xFF)
CD = RGBColor(0x1B, 0x1B, 0x1B)
CG = RGBColor(0x66, 0x66, 0x66)
CL = RGBColor(0xF0, 0xF2, 0xF5)
CB = RGBColor(0xE0, 0xE4, 0xE8)
CN = RGBColor(0x06, 0x2A, 0x66)
FT = "Microsoft YaHei"

def B(slide, c):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = c

def R(slide, l, t, w, h, fill=None, line=None, lw=None, rad=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    if lw: s.line.width = lw
    else: s.line.fill.background()
    return s

def O(slide, l, t, sz, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def T(slide, l, t, w, h, text, size=18, color=None, bold=False, align=PP_ALIGN.LEFT, font=FT):
    if color is None: color = CD
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def H(slide, title, sub=None):
    R(slide, Inches(0.8), Inches(0.5), Inches(0.06), Inches(0.7), fill=C3)
    T(slide, Inches(1.1), Inches(0.35), Inches(11), Inches(0.8), title, size=34, color=C1, bold=True)
    if sub: T(slide, Inches(1.1), Inches(1.0), Inches(11), Inches(0.5), sub, size=16, color=CG)

def PN(slide, n): T(slide, Inches(12.2), Inches(7.15), Inches(1), Inches(0.25), str(n), size=10, color=CB, align=PP_ALIGN.RIGHT)

def GR(slide):
    for i in range(6):
        R(slide, 0, prs.slide_height - Inches(1.2 - i*0.2), prs.slide_width, Inches(1.2 - i*0.2),
          fill=RGBColor(int(11+i*3), int(61+i*3), int(145-i*2)))

print("Helpers loaded")

# SLIDE 1: COVER
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, C1)
for x,y,z,c in [(Inches(0.5),Inches(5),Inches(1.8),C2),(Inches(11.5),Inches(0.5),Inches(1.5),C3),(Inches(10),Inches(5.5),Inches(0.8),C2)]:
    O(s, x, y, z, fill=c)
O(s, Inches(8.5), Inches(2), Inches(5.5), fill=CN)
O(s, Inches(9.5), Inches(3), Inches(3.5), fill=C2)
for i in range(15):
    R(s, Inches(-1+i*1.1), Inches(6), Inches(0.8), Inches(2), fill=C2 if i%2==0 else C1)
R(s, Inches(0.8), Inches(1.8), Inches(6.5), Inches(0.05), fill=C3)
T(s, Inches(0.8), Inches(2.1), Inches(8), Inches(1.5), "\u73cd\u7231\u751f\u547d  \u9884\u9632\u6eba\u6c34", size=56, color=CW, bold=True)
T(s, Inches(0.8), Inches(3.7), Inches(8), Inches(0.8), "\u521d\u4e2d\u751f\u6691\u5047\u5b89\u5168\u6559\u80b2", size=28, color=C3)
R(s, Inches(0.8), Inches(4.6), Inches(2.5), Inches(0.04), fill=C3)
T(s, Inches(0.8), Inches(4.9), Inches(8), Inches(0.6), "\u5b89\u5168\u7b2c\u4e00 \u00b7 \u73cd\u7231\u751f\u547d \u00b7 \u8fdc\u79bb\u5371\u9669\u6c34\u57df", size=16, color=RGBColor(0xAA,0xCC,0xEE))
T(s, Inches(0.8), Inches(6.9), Inches(4), Inches(0.25), "\u4e3b\u9898\u73ed\u4f1a | 2026\u5e74\u6691\u5047", size=12, color=CG)
print("Slide 1")

# SLIDE 2: DATA
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CW)
H(s, "\u6eba\u6c34\u2014\u2014\u4e0d\u5bb9\u5ffd\u89c6\u7684\u5b89\u5168\u5a01\u80c1", "\u6570\u636e\u544a\u8bc9\u4f60\uff0c\u6eba\u6c34\u79bb\u6211\u4eec\u5e76\u4e0d\u9065\u8fdc"); PN(s, 2)
for i,(big,unit,desc,pct,clr) in enumerate([
    ("5.7\u4e07","\u4eba/\u5e74","\u5168\u56fd\u6bcf\u5e74\u6eba\u6c34\u6b7b\u4ea1\u4eba\u6570",0.85,C4),
    ("60%+","\u4e2d\u5c0f\u5b66\u751f","\u6eba\u6c34\u6b7b\u4ea1\u4e2d\u5c0f\u5b66\u751f\u5360\u6bd4",0.60,C5),
    ("7-8\u6708","\u6691\u5047\u671f\u95f4","\u6eba\u6c34\u4e8b\u6545\u6700\u9ad8\u53d1\u65f6\u6bb5",0.90,C1),
]):
    x=Inches(0.8)+i*Inches(4.1); y=Inches(1.8)
    R(s,x,y,Inches(3.7),Inches(3.2),fill=CL,rad=Inches(0.15))
    T(s,x+Inches(0.3),y+Inches(0.3),Inches(3.1),Inches(1.2),big,size=48,color=clr,bold=True)
    T(s,x+Inches(0.3),y+Inches(1.4),Inches(3.1),Inches(0.4),unit,size=18,color=CD)
    R(s,x+Inches(0.3),y+Inches(2.2),Inches(3.1),Inches(0.1),fill=CB,rad=Inches(0.05))
    R(s,x+Inches(0.3),y+Inches(2.2),Inches(3.1*pct),Inches(0.1),fill=clr,rad=Inches(0.05))
    T(s,x+Inches(0.3),y+Inches(2.5),Inches(3.1),Inches(0.5),desc,size=14,color=CG)
R(s,Inches(0.8),Inches(5.5),Inches(11.7),Inches(1.2),fill=RGBColor(0xE8,0xEE,0xF7),rad=Inches(0.1))
R(s,Inches(0.8),Inches(5.5),Inches(0.06),Inches(1.2),fill=C1)
tb=s.shapes.add_textbox(Inches(1.2),Inches(5.6),Inches(10.8),Inches(1));tf=tb.text_frame;tf.word_wrap=True
p=tf.paragraphs[0];p.text="\u636e\u56fd\u5bb6\u536b\u5065\u59d4\u7edf\u8ba1\uff0c\u6eba\u6c34\u662f\u6211\u56fd 1-14 \u5c81\u513f\u7ae5\u610f\u5916\u6b7b\u4ea1\u7684\u7b2c\u4e00\u4f4d\u539f\u56e0\u3002";p.font.size=Pt(18);p.font.color.rgb=C1;p.font.bold=True;p.font.name=FT
p2=tf.add_paragraph();p2.text="\u6bcf\u4e00\u8d77\u6eba\u6c34\u4e8b\u6545\u80cc\u540e\uff0c\u90fd\u662f\u4e00\u4e2a\u5bb6\u5ead\u7684\u7834\u788e\u3002\u9632\u6eba\u6c34\u5b89\u5168\u6559\u80b2\uff0c\u523b\u4e0d\u5bb9\u7f13\uff01";p2.font.size=Pt(15);p2.font.color.rgb=CG;p2.font.name=FT
print("Slide 2")

# SLIDE 3: DANGER ZONES
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CW)
H(s, "\u5371\u9669\u6c34\u57df\u8bc6\u522b", "\u6c34\u9762\u4e4b\u4e0b\uff0c\u5371\u673a\u56db\u4f0f"); PN(s, 3)
for i,(name,danger,detail) in enumerate([
    ("\u91ce\u5916\u6cb3\u6d41","\u6c34\u6d41\u6e4d\u6025\uff0c\u5730\u5f62\u590d\u6742","\u6697\u6d41\u6f29\u6da1\u591a\uff0c\u6c34\u6e29\u6e29\u5dee\u5927\uff0c\u6c34\u4e0b\u4e0d\u53ef\u89c1"),
    ("\u6c34\u5e93\u00b7\u6e56\u6cca","\u6c34\u6df1\u96be\u6d4b\uff0c\u6c34\u6e29\u4f4e","\u5bb9\u6613\u62bd\u7b4b\uff0c\u6de4\u6ce5\u677e\u8f6f\u9677\u4eba\uff0c\u65e0\u6551\u63f4\u8bbe\u65bd"),
    ("\u6d77\u8fb9\u00b7\u6c5f\u8fb9","\u6697\u6d41\u79bb\u5cb8\u6d41","\u6f6e\u6c50\u53d8\u5316\u5feb\uff0c\u65e0\u6551\u751f\u8bbe\u65bd\uff0c\u6551\u63f4\u56f0\u96be"),
    ("\u5de5\u5730\u6c34\u5751","\u96e8\u540e\u79ef\u6c34\uff0c\u8fb9\u6cbf\u677e\u8f6f","\u4f4d\u7f6e\u9690\u853d\u65e0\u4eba\u6ce8\u610f\uff0c\u6c34\u6df1\u4e0d\u53ef\u77e5"),
    ("\u6c60\u5858\u00b7\u6c9f\u6e20","\u6742\u8349\u4e1b\u751f\uff0c\u6c34\u8d28\u6d51\u6d4a","\u6c34\u8349\u4e1b\u751f\u6613\u7f20\u811a\uff0c\u6df1\u6d45\u7a81\u53d8\u5bb9\u6613\u6ed1\u843d"),
]):
    col,row=i%3,i//3
    x=Inches(0.8)+col*Inches(4.1); y=Inches(1.7)+row*Inches(2.7)
    R(s,x,y,Inches(3.8),Inches(2.3),fill=CW,line=CB,lw=Pt(1),rad=Inches(0.12))
    R(s,x,y,Inches(3.8),Inches(0.05),fill=C4,rad=Inches(0.12))
    R(s,x,y+Inches(0.05),Inches(3.8),Inches(0.05),fill=CW)
    T(s,x+Inches(0.25),y+Inches(0.25),Inches(3.3),Inches(0.5),name,size=20,color=CD,bold=True)
    T(s,x+Inches(0.25),y+Inches(0.7),Inches(3.3),Inches(0.35),danger,size=15,color=C4,bold=True)
    T(s,x+Inches(0.25),y+Inches(1.2),Inches(3.3),Inches(0.8),detail,size=13,color=CG)
R(s,Inches(0.8),Inches(6.9),Inches(11.7),Inches(0.35),fill=RGBColor(0xFF,0xF3,0xE0),rad=Inches(0.08))
T(s,Inches(1.1),Inches(6.91),Inches(11),Inches(0.35),"\u770b\u4f3c\u5e73\u9759\u7684\u6c34\u9762\u4e0b\uff0c\u53ef\u80fd\u6697\u85cf\u6f29\u6da1\u3001\u6697\u6d41\u3001\u6c34\u8349\u548c\u6de4\u6ce5\u2014\u2014\u4e0d\u662f\u6240\u6709\u5371\u9669\u90fd\u8089\u773c\u53ef\u89c1",size=13,color=C5,align=PP_ALIGN.CENTER)
print("Slide 3")

# SLIDE 4: SIX DONTS
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CN); GR(s)
R(s,Inches(0.8),Inches(0.5),Inches(0.06),Inches(0.7),fill=C3)
T(s,Inches(1.1),Inches(0.35),Inches(11),Inches(0.8),"\u9632\u6eba\u6c34 \u00b7 \u516d\u4e0d\u51c6",size=38,color=CW,bold=True)
T(s,Inches(1.1),Inches(1.0),Inches(11),Inches(0.5),"\u4e0d\u662f\u300c\u5efa\u8bae\u300d\uff0c\u662f\u300c\u5fc5\u987b\u300d\u2014\u2014\u6bcf\u4e00\u6761\u90fd\u662f\u7528\u751f\u547d\u6362\u6765\u7684\u6559\u8bad",size=16,color=RGBColor(0x90,0xBB,0xEE))
PN(s,4)
for i,(num,rule) in enumerate([
    (1,"\u4e0d\u51c6\u79c1\u81ea\u4e0b\u6c34\u6e38\u6cf3"),(2,"\u4e0d\u51c6\u64c5\u81ea\u4e0e\u4ed6\u4eba\u7ed3\u4f34\u6e38\u6cf3"),
    (3,"\u4e0d\u51c6\u5728\u65e0\u5bb6\u957f\u6216\u8001\u5e08\u5e26\u9886\u7684\u60c5\u51b5\u4e0b\u6e38\u6cf3"),(4,"\u4e0d\u51c6\u5230\u65e0\u5b89\u5168\u8bbe\u65bd\u3001\u65e0\u6551\u62a4\u4eba\u5458\u7684\u6c34\u57df\u6e38\u6cf3"),
    (5,"\u4e0d\u51c6\u5230\u4e0d\u719f\u6089\u7684\u6c34\u57df\u6e38\u6cf3"),(6,"\u4e0d\u51c6\u4e0d\u719f\u6089\u6c34\u6027\u7684\u5b66\u751f\u64c5\u81ea\u4e0b\u6c34\u65bd\u6551"),
]):
    row,col=i//2,i%2
    x=Inches(0.8)+col*Inches(6.2); y=Inches(1.8)+row*Inches(1.65)
    R(s,x,y,Inches(5.8),Inches(1.35),fill=RGBColor(0x0E,0x49,0x9B),line=RGBColor(0x1A,0x60,0xB8),lw=Pt(1),rad=Inches(0.1))
    O(s,x+Inches(0.25),y+Inches(0.25),Inches(0.75),fill=C3)
    T(s,x+Inches(0.25),y+Inches(0.32),Inches(0.75),Inches(0.6),"0"+str(num),size=22,color=CW,bold=True,align=PP_ALIGN.CENTER)
    T(s,x+Inches(1.2),y+Inches(0.3),Inches(4.3),Inches(0.7),rule,size=20,color=CW)
R(s,Inches(2),Inches(6.85),Inches(9.3),Inches(0.35),fill=C4,rad=Inches(0.18))
T(s,Inches(2.3),Inches(6.85),Inches(8.7),Inches(0.35),"\u4ee5\u4e0a\u516d\u6761\uff0c\u8fdd\u53cd\u4efb\u4f55\u4e00\u6761\u90fd\u53ef\u80fd\u4ed8\u51fa\u751f\u547d\u7684\u4ee3\u4ef7",size=14,color=CW,bold=True,align=PP_ALIGN.CENTER)
print("Slide 4")

# SLIDE 5: RESCUE OTHERS
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CW)
H(s, "\u53d1\u73b0\u6709\u4eba\u6eba\u6c34\uff0c\u5e94\u8be5\u600e\u4e48\u505a\uff1f", "\u6b63\u786e\u7684\u53cd\u5e94\u80fd\u633d\u6551\u751f\u547d\u2014\u2014\u51b7\u9759\u6bd4\u52c7\u6562\u66f4\u91cd\u8981"); PN(s,5)
for i,(num,title,desc,clr) in enumerate([
    ("1","\u5927\u58f0\u547c\u6551","\u7acb\u5373\u9ad8\u58f0\u547c\u557c\u5f15\u8d77\u6ce8\u610f\n\u62e8\u6253 110 \u62a5\u8b66\u3001120 \u6025\u6551",C1),
    ("2","\u5bfb\u627e\u5de5\u5177","\u5bfb\u627e\u6551\u751f\u5708\u3001\u7af9\u7aff\u3001\u7ef3\u5b50\n\u629b\u7ed9\u6eba\u6c34\u8005\uff0c\u8ba9\u5176\u6293\u4f4f",C2),
    ("3","\u5207\u52ff\u4e0b\u6c34","\u672a\u6210\u5e74\u4eba\u7edd\u4e0d\u8981\u4e0b\u6c34\u65bd\u6551\uff01\n\u8d38\u7136\u4e0b\u6c34\u6781\u5176\u5371\u9669",C4),
    ("4","\u534f\u52a9\u6551\u63f4","\u914d\u5408\u6210\u5e74\u4eba\u8fdb\u884c\u6551\u63f4\n\u4e0a\u5cb8\u540e\u534f\u52a9\u4fdd\u6696\u548c\u6025\u6551",C3),
]):
    x=Inches(0.8)+i*Inches(3.15); y=Inches(1.7)
    R(s,x,y,Inches(2.85),Inches(4.5),fill=CW,line=CB,lw=Pt(0.5),rad=Inches(0.12))
    R(s,x,y,Inches(2.85),Inches(1.5),fill=clr,rad=Inches(0.12))
    R(s,x,y+Inches(1.2),Inches(2.85),Inches(0.3),fill=clr)
    T(s,x+Inches(0.3),y+Inches(0.15),Inches(2.25),Inches(0.7),num,size=48,color=CW,bold=True,font="Arial")
    T(s,x+Inches(0.3),y+Inches(0.75),Inches(2.25),Inches(0.5),title,size=20,color=CW,bold=True)
    T(s,x+Inches(0.25),y+Inches(1.7),Inches(2.35),Inches(2.5),desc,size=14,color=CD)
R(s,Inches(0.8),Inches(6.6),Inches(11.7),Inches(0.6),fill=RGBColor(0xFD,0xED,0xEC),line=C4,lw=Pt(1.5),rad=Inches(0.1))
T(s,Inches(1.2),Inches(6.63),Inches(10.9),Inches(0.55),"\u5207\u8bb0\uff1a\u4fdd\u5df1\u6551\u4eba\uff01\u4fdd\u8bc1\u81ea\u8eab\u5b89\u5168\u662f\u7b2c\u4e00\u539f\u5219\uff0c\u4e0d\u4e0b\u6c34\u65bd\u6551\u4e0d\u662f\u51b7\u6f20\uff0c\u662f\u7406\u667a\u3002",size=15,color=C4,bold=True,align=PP_ALIGN.CENTER)
print("Slide 5")

# SLIDE 6: SELF RESCUE
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CW)
H(s, "\u4e0d\u614e\u843d\u6c34\uff0c\u5982\u4f55\u81ea\u6551\uff1f", "\u614c\u4e71\u662f\u6700\u5927\u7684\u654c\u4eba\u2014\u2014\u8bb0\u4f4f\u8fd9 6 \u4e2a\u5173\u952e\u52a8\u4f5c"); PN(s,6)
for i,(title,desc) in enumerate([
    ("\u4fdd\u6301\u51b7\u9759","\u4e0d\u8981\u614c\u5f20\u6323\u624e\uff0c\u614c\u4e71\u4f1a\n\u52a0\u901f\u4f53\u529b\u6d88\u8017\uff0c\u6076\u5316\u60c5\u51b5"),
    ("\u4ef0\u9762\u6f02\u6d6e","\u8eab\u4f53\u540e\u4ef0\u53e3\u9f3b\u9732\u51fa\u6c34\u9762\n\u6df1\u5438\u6d45\u547c\uff0c\u4fdd\u5b58\u4f53\u529b"),
    ("\u4e3e\u624b\u547c\u6551","\u5355\u81c2\u6325\u52a8\u6c42\u6551\u7b49\u5f85\u6551\u63f4\n\u4e0d\u8981\u62fc\u547d\u5411\u5cb8\u8fb9\u6e38"),
    ("\u8131\u6389\u91cd\u7269","\u8131\u6389\u978b\u5b50\u548c\u539a\u91cd\u5916\u5957\n\u4fdd\u7559\u90e8\u5206\u8863\u7269\u9632\u6b62\u5931\u6e29"),
    ("\u501f\u52a9\u6d6e\u7269","\u6293\u4f4f\u6728\u677f\u7a7a\u6c34\u74f6\u7b49\u6f02\u6d6e\u7269\n\u4efb\u4f55\u80fd\u6d6e\u7684\u4e1c\u897f\u90fd\u80fd\u6551\u547d"),
    ("\u62bd\u7b4b\u81ea\u6551","\u624b\u6307\u62bd\u7b4b\uff1a\u63e1\u62f3\u2192\u7528\u529b\u5f20\u5f00\n\u5c0f\u817f\u62bd\u7b4b\uff1a\u8e6c\u76f4\u817f\u2192\u52fe\u811a\u5c16"),
]):
    col,row=i%3,i//3
    x=Inches(0.8)+col*Inches(4.1); y=Inches(1.7)+row*Inches(2.6)
    R(s,x,y,Inches(3.8),Inches(2.25),fill=CW,line=CB,lw=Pt(0.5),rad=Inches(0.1))
    R(s,x,y,Inches(0.06),Inches(2.25),fill=C2,rad=Inches(0.03))
    T(s,x+Inches(0.3),y+Inches(0.2),Inches(3.2),Inches(0.5),title,size=20,color=C1,bold=True)
    T(s,x+Inches(0.3),y+Inches(0.7),Inches(3.2),Inches(1.2),desc,size=14,color=CD)
print("Slide 6")

# SLIDE 7: EMERGENCY NUMBERS
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CN); GR(s)
R(s,Inches(0.8),Inches(0.5),Inches(0.06),Inches(0.7),fill=C3)
T(s,Inches(1.1),Inches(0.35),Inches(11),Inches(0.8),"\u7d27\u6025\u6c42\u52a9\u7535\u8bdd",size=38,color=CW,bold=True)
T(s,Inches(1.1),Inches(1.0),Inches(11),Inches(0.5),"\u8bb0\u4e0b\u8fd9\u4e9b\u53f7\u7801\u2014\u2014\u5173\u952e\u65f6\u523b\u80fd\u6551\u547d",size=16,color=RGBColor(0x90,0xBB,0xEE))
PN(s,7)
for i,(num,label,desc) in enumerate([
    ("110","\u62a5\u8b66\u7535\u8bdd","\u9047\u5230\u5371\u9669\u60c5\u51b5\n\u7b2c\u4e00\u65f6\u95f4\u62e8\u6253"),
    ("120","\u533b\u7597\u6025\u6551","\u6709\u4eba\u53d7\u4f24\u6216\u6eba\u6c34\n\u9700\u8981\u7d27\u6025\u533b\u7597\u6551\u52a9"),
    ("119","\u6d88\u9632\u6551\u63f4","\u706b\u707e\u3001\u88ab\u56f0\u7b49\n\u7d27\u6025\u6551\u63f4\u60c5\u51b5"),
    ("12395","\u6c34\u4e0a\u6551\u63f4","\u56fd\u5bb6\u6c34\u4e0a\u9047\u9669\n\u6c42\u6551\u4e13\u7ebf\u7535\u8bdd"),
]):
    x=Inches(0.8)+i*Inches(3.15); y=Inches(1.8)
    R(s,x,y,Inches(2.85),Inches(4.2),fill=RGBColor(0x0E,0x49,0x9B),line=RGBColor(0x1A,0x60,0xB8),lw=Pt(1),rad=Inches(0.12))
    O(s,x+Inches(0.75),y+Inches(0.4),Inches(1.2),fill=C3)
    T(s,x+Inches(0.75),y+Inches(0.55),Inches(1.2),Inches(0.9),num,size=28,color=CW,bold=True,align=PP_ALIGN.CENTER,font="Arial")
    T(s,x+Inches(0.2),y+Inches(1.8),Inches(2.45),Inches(0.5),label,size=18,color=CW,bold=True,align=PP_ALIGN.CENTER)
    T(s,x+Inches(0.2),y+Inches(2.4),Inches(2.45),Inches(1.2),desc,size=13,color=RGBColor(0xBB,0xCC,0xEE),align=PP_ALIGN.CENTER)
R(s,Inches(0.8),Inches(6.3),Inches(11.7),Inches(0.9),fill=RGBColor(0x0E,0x49,0x9B),line=RGBColor(0x1A,0x60,0xB8),lw=Pt(1),rad=Inches(0.08))
T(s,Inches(1.1),Inches(6.35),Inches(11),Inches(0.8),"\u6253\u7535\u8bdd\u65f6\u8bf7\u8bf4\u6e05\uff1a\u2460 \u53d1\u751f\u4e86\u4ec0\u4e48  ->  \u2461 \u5177\u4f53\u4f4d\u7f6e\uff08\u9644\u8fd1\u9192\u76ee\u6807\u5fd7\uff09 ->  \u2462 \u6eba\u6c34\u8005\u60c5\u51b5  ->  \u2463 \u4fdd\u6301\u901a\u8bdd\uff0c\u6d3e\u4eba\u63a5\u5e94",size=14,color=CW,align=PP_ALIGN.CENTER)
print("Slide 7")

# SLIDE 8: PARENT GUIDE
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CW)
H(s, "\u7ed9\u5bb6\u957f\u7684\u63d0\u9192", "\u5b69\u5b50\u5b89\u5168\uff0c\u5bb6\u957f\u662f\u7b2c\u4e00\u8d23\u4efb\u4eba"); PN(s,8)
for i,(title,desc) in enumerate([("\u77e5\u9053\u53bb\u5411","\u5b69\u5b50\u53bb\u4e86\u54ea\u91cc"),("\u77e5\u9053\u540c\u4f34","\u548c\u8c01\u5728\u4e00\u8d77"),("\u77e5\u9053\u5185\u5bb9","\u53bb\u505a\u4ec0\u4e48\u6d3b\u52a8"),("\u77e5\u9053\u5f52\u65f6","\u4ec0\u4e48\u65f6\u5019\u56de\u6765")]):
    x=Inches(0.8)+i*Inches(3.15); y=Inches(1.7)
    R(s,x,y,Inches(2.85),Inches(2.5),fill=CW,line=CB,lw=Pt(1),rad=Inches(0.12))
    R(s,x,y,Inches(2.85),Inches(0.7),fill=C3,rad=Inches(0.12))
    R(s,x,y+Inches(0.5),Inches(2.85),Inches(0.2),fill=C3)
    T(s,x+Inches(0.2),y+Inches(0.08),Inches(2.45),Inches(0.55),title,size=22,color=CW,bold=True,align=PP_ALIGN.CENTER)
    T(s,x+Inches(0.2),y+Inches(1.0),Inches(2.45),Inches(1.0),desc,size=16,color=CD,align=PP_ALIGN.CENTER)
R(s,Inches(0.8),Inches(4.6),Inches(11.7),Inches(2.5),fill=CL,rad=Inches(0.1))
tb=s.shapes.add_textbox(Inches(1.2),Inches(4.8),Inches(10.8),Inches(2.1));tf=tb.text_frame;tf.word_wrap=True
lines=[("\u5bb6\u957f\u5b89\u5168\u6307\u5357",20,C1,True),("",8,CD,False),("\u2022 \u5e26\u5b69\u5b50\u53bb\u6b63\u89c4\u6e38\u6cf3\u9986\uff0c\u5168\u7a0b\u966a\u540c\u76d1\u62a4\uff0c\u52ff\u8ba9\u5b69\u5b50\u79bb\u5f00\u89c6\u7ebf",15,CD,False),("\u2022 \u5bb6\u4e2d\u6c34\u7f38\u3001\u6c34\u6876\u53ca\u65f6\u52a0\u76d6\uff1b\u536b\u751f\u95f4\u4e0d\u7528\u65f6\u9501\u95e8",15,CD,False),("\u2022 \u6559\u4f1a\u5b69\u5b50\u57fa\u672c\u7684\u81ea\u6551\u548c\u6c42\u6551\u65b9\u6cd5\uff0c\u719f\u8bb0\u7d27\u6025\u7535\u8bdd",15,CD,False),("\u2022 \u5916\u51fa\u6e38\u73a9\u65f6\u65f6\u523b\u5173\u6ce8\u5b69\u5b50\u52a8\u5411\uff0c\u4e0d\u5728\u6c34\u8fb9\u73a9\u624b\u673a",15,CD,False),("\u2022 \u53d1\u73b0\u5176\u4ed6\u5b69\u5b50\u72ec\u81ea\u5728\u6c34\u8fb9\u73a9\u800d\uff0c\u4e3b\u52a8\u63d0\u9192\u6216\u8054\u7cfb\u5bb6\u957f",15,CD,False)]
for i,(txt,sz,clr,bld) in enumerate(lines):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=clr;p.font.bold=bld;p.font.name=FT
print("Slide 8")

# SLIDE 9: PLEDGE
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, C1)
for x,y,z,c in [(Inches(-0.5),Inches(5.5),Inches(3),C2),(Inches(11),Inches(-0.5),Inches(2),C3),(Inches(10),Inches(6),Inches(1.2),C2)]:
    O(s,x,y,z,fill=c)
T(s,Inches(1.5),Inches(0.8),Inches(10),Inches(1.2),"\u6211\u627f\u8bfa",size=52,color=CW,bold=True,align=PP_ALIGN.CENTER)
pledge=["\u73cd\u7231\u751f\u547d\uff0c\u8fdc\u79bb\u5371\u9669\u6c34\u57df","\u4e0d\u79c1\u81ea\u4e0b\u6c34\u6e38\u6cf3\uff0c\u4e0d\u64c5\u81ea\u7ed3\u4f34\u6e38\u6cf3","\u4e0d\u5728\u65e0\u5bb6\u957f\u5e26\u9886\u65f6\u6e38\u6cf3","\u4e0d\u5230\u65e0\u5b89\u5168\u8bbe\u65bd\u3001\u65e0\u6551\u63f4\u4eba\u5458\u7684\u6c34\u57df\u6e38\u6cf3","","\u5b66\u4f1a\u57fa\u672c\u7684\u81ea\u62a4\u81ea\u6551\u65b9\u6cd5","\u53d1\u73b0\u9669\u60c5\u53ca\u65f6\u62a5\u544a\uff0c\u4e0d\u76f2\u76ee\u4e0b\u6c34\u65bd\u6551","","\u8ba9\u5b89\u5168\u966a\u4f34\u6574\u4e2a\u6691\u5047"]
tb=s.shapes.add_textbox(Inches(2.5),Inches(2.2),Inches(8),Inches(4));tf=tb.text_frame;tf.word_wrap=True
for i,line in enumerate(pledge):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    if not line: p.text="";p.space_after=Pt(8)
    else: p.text="  "+line;p.font.size=Pt(20);p.font.color.rgb=CW;p.font.name=FT;p.space_after=Pt(6);p.alignment=PP_ALIGN.CENTER
R(s,Inches(5),Inches(6.2),Inches(3.3),Inches(0.04),fill=C3)
T(s,Inches(2),Inches(6.4),Inches(9.3),Inches(0.5),"\u795d\u540c\u5b66\u4eec\u5ea6\u8fc7\u4e00\u4e2a\u5e73\u5b89\u3001\u5065\u5eb7\u3001\u5feb\u4e50\u7684\u6691\u5047\uff01",size=16,color=RGBColor(0xAA,0xCC,0xEE),align=PP_ALIGN.CENTER)
PN(s,9)
print("Slide 9")

# SLIDE 10: END
s = prs.slides.add_slide(prs.slide_layouts[6]); B(s, CN); GR(s)
T(s,Inches(1),Inches(2),Inches(11),Inches(1.5),"\u8c22\u8c22\u89c2\u770b",size=52,color=CW,bold=True,align=PP_ALIGN.CENTER)
R(s,Inches(5.2),Inches(3.5),Inches(2.9),Inches(0.04),fill=C3)
T(s,Inches(1),Inches(3.8),Inches(11),Inches(0.8),"\u5b89\u5168\u65e0\u5c0f\u4e8b  \u9632\u60a3\u4e8e\u672a\u7136",size=24,color=C3,align=PP_ALIGN.CENTER)
T(s,Inches(1),Inches(5),Inches(11),Inches(0.6),"2026\u5e74\u6691\u5047\u9632\u6eba\u6c34\u5b89\u5168\u6559\u80b2\u4e3b\u9898\u73ed\u4f1a",size=16,color=RGBColor(0x90,0xBB,0xEE),align=PP_ALIGN.CENTER)
PN(s,10)
print("Slide 10")

# SAVE
output = "C:/Users/89728/Documents/Codex/2026-06-06/git/test-git-link/暑假防溺水安全教育.pptx"
prs.save(output)
print("SAVED: " + output)
print("Total slides: " + str(len(prs.slides)))
