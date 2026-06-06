# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DEEP_BLUE = RGBColor(0x0D, 0x47, 0xA1)
LIGHT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x23, 0x2E)
WARM_RED = RGBColor(0xE5, 0x39, 0x35)
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)
LIGHT_BG = RGBColor(0xE3, 0xF2, 0xFD)

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    if color is None: color = DARK_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    if font_name: p.font.name = font_name
    p.alignment = alignment

def add_wave(slide, color=None):
    if color is None: color = LIGHT_BLUE
    shape = slide.shapes.add_shape(MSO_SHAPE.WAVE, 0, prs.slide_height - Inches(1.5), prs.slide_width, Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ==== SLIDE 1: TITLE ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_BLUE)
for x, y, sz, c in [(Inches(10), Inches(1), Inches(3), RGBColor(0x15, 0x65, 0xC0)), (Inches(11), Inches(5), Inches(2), RGBColor(0x19, 0x76, 0xD2)), (Inches(1), Inches(5.5), Inches(1.5), LIGHT_BLUE)]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    circle.fill.solid(); circle.fill.fore_color.rgb = c; circle.line.fill.background()
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(2), "珍爱生命  预防溺水", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1), "初中生暑假安全教育主题班会", font_size=28, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.8), "安全第一 | 珍爱生命 | 远离危险水域", font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# ==== SLIDE 2: STATISTICS ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG); add_wave(slide)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1), "溺水——青少年意外死亡的头号杀手", font_size=36, color=DEEP_BLUE, bold=True, font_name="Microsoft YaHei")
stats = [("每年约5.7万人", "死于溺水", WARM_RED), ("其中中小学生", "占比超过60%", ACCENT_ORANGE), ("暑假7-8月", "事故高发期", DEEP_BLUE)]
for i, (num, desc, c) in enumerate(stats):
    x = Inches(1) + i * Inches(3.8)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), Inches(2), Inches(2.2), Inches(2.2))
    circle.fill.solid(); circle.fill.fore_color.rgb = WHITE; circle.line.color.rgb = c; circle.line.width = Pt(3)
    add_textbox(slide, x + Inches(0.5), Inches(2.3), Inches(2.2), Inches(0.8), num, font_size=28, color=c, bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
    add_textbox(slide, x + Inches(0.5), Inches(3.1), Inches(2.2), Inches(0.6), desc, font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
add_textbox(slide, Inches(1), Inches(5), Inches(11), Inches(1.5), "据国家卫健委统计，溺水是我国1-14岁儿童意外死亡的第一位原因。\n每一起溺水事故的背后，都是一个家庭的破碎。\n防溺水安全教育，刻不容缓！", font_size=18, color=DARK_TEXT, font_name="Microsoft YaHei")

# ==== SLIDE 3: DANGEROUS AREAS ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG); add_wave(slide)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(1), "⚠️ 这些地方很危险！", font_size=36, color=WARM_RED, bold=True, font_name="Microsoft YaHei")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5), "看似平静的水面下，可能暗藏漩涡、暗流、水草和淤泥", font_size=18, color=DARK_TEXT, font_name="Microsoft YaHei")
areas = [("野外河流", "水流湍急，水下地形复杂\n暗流漩涡多，水温温差大"), ("水库/湖泊", "水深不可测，水温低\n容易抽筋，淤泥陷人"), ("海边/江边", "潮汐暗流，离岸流\n无救生设施，救援困难"), ("工地水坑", "雨后积水，边沿松软\n隐蔽危险，无人看管"), ("池塘/沟渠", "水质浑浊，暗藏水草\n深浅不一，容易滑落")]
for i, (title, desc) in enumerate(areas):
    col, row = i % 3, i // 3
    x = Inches(0.8) + col * Inches(4.1); y = Inches(1.8) + row * Inches(2.6)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.2))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = RGBColor(0xFF, 0xCD, 0xD2); card.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.6), title, font_size=22, color=WARM_RED, bold=True, font_name="Microsoft YaHei")
    add_textbox(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.2), desc, font_size=15, color=DARK_TEXT, font_name="Microsoft YaHei")

# ==== SLIDE 4: SIX DON\'TS ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(1), "防溺水六不准", font_size=40, color=WHITE, bold=True, font_name="Microsoft YaHei")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5), "每一位同学都必须牢记于心！", font_size=20, color=LIGHT_BLUE, font_name="Microsoft YaHei")
rules = [("1", "不准私自下水游泳"), ("2", "不准擅自与他人结伴游泳"), ("3", "不准在无家长或老师带领的情况下游泳"), ("4", "不准到无安全设施、无救护人员的水域游泳"), ("5", "不准到不熟悉的水域游泳"), ("6", "不准不熟悉水性的学生擅自下水施救")]
for i, (num, rule) in enumerate(rules):
    y = Inches(1.8) + i * Inches(0.85)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), y + Pt(2), Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = LIGHT_BLUE; circle.line.fill.background()
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(22); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    add_textbox(slide, Inches(2.3), y, Inches(9), Inches(0.6), rule, font_size=22, color=WHITE, font_name="Microsoft YaHei")

# ==== SLIDE 5: WHAT IF SOMEONE DROWNS ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG); add_wave(slide)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(1), "发现有人溺水怎么办？", font_size=36, color=DEEP_BLUE, bold=True, font_name="Microsoft YaHei")
steps = [("第一步：大声呼救", "立即大声呼喊，引起周围成年人注意\n拨打110报警、120急救电话"), ("第二步：寻找救援工具", "寻找救生圈、竹竿、绳子等\n抛给溺水者，让其抓住"), ("第三步：绝不盲目下水", "未成年人不要下水施救！\n你不是专业救生员，贸然下水极其危险"), ("第四步：协助救援", "配合成年人进行救援\n救上岸后协助进行简单急救")]
for i, (title, desc) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3.15); y = Inches(1.6)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(3.8))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = DEEP_BLUE; card.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.5), Inches(0.6), title, font_size=17, color=DEEP_BLUE, bold=True, font_name="Microsoft YaHei")
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.5), Inches(2.5), desc, font_size=14, color=DARK_TEXT, font_name="Microsoft YaHei")
highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.8), Inches(11.3), Inches(0.9))
highlight.fill.solid(); highlight.fill.fore_color.rgb = WARM_RED; highlight.line.fill.background()
add_textbox(slide, Inches(1.3), Inches(5.85), Inches(10.7), Inches(0.8), "⚠️ 切记：保己救人！保证自身安全是第一位的，不要做无谓的牺牲！", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# ==== SLIDE 6: SELF-RESCUE ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG); add_wave(slide)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(1), "自己不慎落水如何自救？", font_size=36, color=DEEP_BLUE, bold=True, font_name="Microsoft YaHei")
self_rescue = [("保持冷静", "不要慌张，不要胡乱挣扎\n慌乱会加速体力消耗"), ("仰面漂浮", "身体后仰，口鼻露出水面\n深吸浅呼，保存体力"), ("呼救待援", "挥手呼救，等待救援\n不要拼命游向岸边"), ("去除重物", "脱掉鞋子、重衣服\n但不要脱掉所有衣物防失温"), ("利用浮物", "抓住身边的漂浮物\n木板、空水瓶都能救命"), ("抽筋处理", "手指抽筋：握拳用力张开\n小腿抽筋：蹬直腿勾脚尖")]
for i, (title, desc) in enumerate(self_rescue):
    col, row = i % 3, i // 3
    x = Inches(0.8) + col * Inches(4.1); y = Inches(1.5) + row * Inches(2.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.4))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = LIGHT_BLUE; card.line.width = Pt(1)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5), title, font_size=20, color=DEEP_BLUE, bold=True, font_name="Microsoft YaHei")
    add_textbox(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.4), desc, font_size=15, color=DARK_TEXT, font_name="Microsoft YaHei")

# ==== SLIDE 7: EMERGENCY NUMBERS ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(1), "牢记紧急电话", font_size=40, color=WHITE, bold=True, font_name="Microsoft YaHei")
phones = [("110", "报警电话"), ("120", "急救电话"), ("119", "火警电话"), ("12395", "水上救援")]
for i, (num, label) in enumerate(phones):
    x = Inches(2) + i * Inches(2.5); y = Inches(1.8)
    add_textbox(slide, x, y, Inches(2), Inches(0.9), num, font_size=44, color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Arial")
    add_textbox(slide, x, y + Inches(1), Inches(2), Inches(0.5), label, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(2.5), "‼️ 打电话时要说清楚：\n① 发生了什么（有人溺水）\n② 具体位置（附近有什么标志性建筑）\n③ 溺水者情况（人数、状态）\n④ 保持电话畅通，派人接应", font_size=18, color=WHITE, font_name="Microsoft YaHei")

# ==== SLIDE 8: PARENT GUIDE ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG); add_wave(slide)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(1), "给家长的四知道", font_size=36, color=DEEP_BLUE, bold=True, font_name="Microsoft YaHei")
parent_tips = [("知道去向", "孩子去了哪里？"), ("知道同伴", "孩子和谁一起？"), ("知道内容", "孩子去做什么？"), ("知道归时", "孩子何时回来？")]
for i, (title, desc) in enumerate(parent_tips):
    x = Inches(1) + i * Inches(3); y = Inches(1.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.6), Inches(2.2))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = ACCENT_ORANGE; card.line.width = Pt(2)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.2), Inches(0.6), title, font_size=22, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.2), Inches(0.8), desc, font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(2), "家长请做到：\n• 带孩子去正规游泳馆，全程陪同监护\n• 家中水缸、水桶加盖，卫生间不用时关门\n• 教会孩子基本的自救和求救方法\n• 外出游玩时时刻关注孩子动向", font_size=18, color=DARK_TEXT, font_name="Microsoft YaHei")

# ==== SLIDE 9: PLEDGE ====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DEEP_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.5), "我承诺", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
pledge = "珍爱生命，远离危险水域\n不私自下水游泳，不擅自结伴游泳\n不在无家长带领时游泳\n不到无安全设施的水域游泳\n\n学会基本的自护自救方法\n发现险情及时报告，不盲目施救\n\n让安全陪伴整个暑假！"
add_textbox(slide, Inches(2), Inches(2.2), Inches(9), Inches(4.5), pledge, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")
add_textbox(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4), "祝同学们度过一个平安、快乐的暑假！", font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER, font_name="Microsoft YaHei")

# ==== SAVE ====
output_path = r"C:\Users\89728\Documents\Codex\2026-06-06\git\test-git-link\暑假防溺水安全教育.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
